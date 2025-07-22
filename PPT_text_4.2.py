import tkinter as tk
from tkinter import filedialog, messagebox
import pandas as pd
import re
import os
from pptx import Presentation

# === GLOBAL LISTS (re‑used each extraction) ===
data_all = []
data_vn  = []

# ──────────────────────────────────────────────
# Helper utilities
# ──────────────────────────────────────────────
def contains_vietnamese(text: str) -> bool:
    """Return True if the string contains at least one Vietnamese‑specific character."""
    if not isinstance(text, str):
        return False
    return bool(re.search(
        r"[àáảãạăằắẳẵặâầấẩẫậèéẻẽẹêềếểễệ"
        r"ìíỉĩịòóỏõọôồốổỗộơờớởỡợ"
        r"ùúủũụưừứửữựỳýỷỹỵđ]",
        text, re.IGNORECASE
    ))

def safe_get_shape_type(shape):
    """Return shape_type or None when python‑pptx cannot identify the type."""
    try:
        return shape.shape_type
    except NotImplementedError:
        return None

# ──────────────────────────────────────────────
# Button #1 – Extract both *_all_text.xlsx & *_vietnamese_only.xlsx
# ──────────────────────────────────────────────
def extract_text_both():
    pptx_path = filedialog.askopenfilename(
        filetypes=[("PowerPoint files", "*.pptx")],
        title="Select PPTX to extract"
    )
    if not pptx_path:
        return

    prs = Presentation(pptx_path)
    data_all.clear()
    data_vn.clear()

    def walk(shape, slide_idx, idx_path=""):
        shape_type = safe_get_shape_type(shape)

        if shape_type == 6:                                  # Group shape
            for k, sub in enumerate(shape.shapes):
                nested = f"{idx_path}.{k}" if idx_path else str(k)
                walk(sub, slide_idx, nested)

        elif shape.has_text_frame:                           # Regular text shape
            text = shape.text

            # All‑text file
            data_all.append({
                "Slide": slide_idx + 1,
                "ShapeIndex": idx_path,
                "OriginalText": text
            })

            # Vietnamese‑only file (same rows, blank if not VN)
            data_vn.append({
                "Slide": slide_idx + 1,
                "ShapeIndex": idx_path,
                "OriginalText": text if contains_vietnamese(text) else "",
                "TranslatedText": ""
            })

    # Traverse slides
    for i, slide in enumerate(prs.slides):
        for j, shp in enumerate(slide.shapes):
            walk(shp, i, str(j))

    base = os.path.splitext(pptx_path)[0]
    out_all = base + "_all_text.xlsx"
    out_vn  = base + "_vietnamese_only.xlsx"

    pd.DataFrame(data_all).to_excel(out_all, index=False)
    pd.DataFrame(data_vn ).to_excel(out_vn , index=False)

    messagebox.showinfo(
        "Extraction complete",
        f"✅ Created:\n• {os.path.basename(out_all)}\n• {os.path.basename(out_vn)}"
    )

# ──────────────────────────────────────────────
# Button #2 – Apply translations back to PPTX
# ──────────────────────────────────────────────
def apply_translation():
    pptx_path = filedialog.askopenfilename(
        filetypes=[("PowerPoint files", "*.pptx")],
        title="Original PPTX to receive translations"
    )
    if not pptx_path:
        return

    excel_path = filedialog.askopenfilename(
        filetypes=[("Excel files", "*.xlsx")],
        title="Excel containing 'TranslatedText' column"
    )
    if not excel_path:
        return

    try:
        prs = Presentation(pptx_path)
        df  = pd.read_excel(excel_path)

        def nested_shape(shapes, path: str):
            idxs = list(map(int, path.split(".")))
            sh   = shapes[idxs[0]]
            for n in idxs[1:]:
                sh = sh.shapes[n]
            return sh

        for _, row in df.iterrows():
            new_txt = row.get("TranslatedText")

            if pd.isna(new_txt) or str(new_txt).strip() == "":
                continue

            slide_idx  = int(row["Slide"]) - 1
            shape_path = str(row["ShapeIndex"]).strip()

            try:
                shp = nested_shape(prs.slides[slide_idx].shapes, shape_path)
                if not shp.has_text_frame:
                    continue
            except Exception as e:
                print(f"⚠️  Skip: slide {slide_idx+1} index '{shape_path}': {e}")
                continue

            tf   = shp.text_frame
            parts = str(new_txt).split("\n")

            for i, txt in enumerate(parts):
                if i < len(tf.paragraphs):
                    p = tf.paragraphs[i]
                    if p.runs:
                        p.runs[0].text = txt
                        for r in p.runs[1:]:
                            r.text = ""
                    else:
                        p.add_run().text = txt
                else:
                    tf.add_paragraph().text = txt

            # remove leftover old paragraphs
            for j in range(len(parts), len(tf.paragraphs)):
                tf.paragraphs[j].clear()

        out_pptx = os.path.splitext(pptx_path)[0] + "_translated.pptx"
        prs.save(out_pptx)
        messagebox.showinfo("Done", f"✅ Saved:\n{out_pptx}")

    except Exception as e:
        messagebox.showerror("Error", str(e))

# ──────────────────────────────────────────────
# Button #3 – Copy 3rd‑column values → replace 4th column & rename header
# ──────────────────────────────────────────────
def replace_excel_column():
    src_path = filedialog.askopenfilename(
        filetypes=[("Excel files", "*.xlsx")],
        title="Source Excel (take 3rd column)"
    )
    if not src_path:
        return

    tgt_path = filedialog.askopenfilename(
        filetypes=[("Excel files", "*.xlsx")],
        title="Target Excel (replace 4th column)"
    )
    if not tgt_path:
        return

    try:
        df_src = pd.read_excel(src_path)
        df_tgt = pd.read_excel(tgt_path)

        if df_src.shape[1] < 3:
            messagebox.showerror("Error", "❌ Source requires ≥3 columns.")
            return
        if df_tgt.shape[1] < 4:
            messagebox.showerror("Error", "❌ Target requires ≥4 columns.")
            return

        # Values (whole column incl. header row)
        col_vals          = df_src.iloc[:, 2].values
        df_tgt.iloc[:, 3] = col_vals

        # Rename header only (4th column)
        df_tgt.columns.values[3] = "TranslatedText"

        out_xlsx = os.path.splitext(tgt_path)[0] + "_modified.xlsx"
        df_tgt.to_excel(out_xlsx, index=False)
        messagebox.showinfo("Success", f"✅ Saved:\n{out_xlsx}")

    except Exception as e:
        messagebox.showerror("Error", str(e))

# ──────────────────────────────────────────────
# Tkinter GUI
# ──────────────────────────────────────────────
root = tk.Tk()
root.title("PPTX & Excel Text Utility")
root.geometry("440x300")
root.configure(bg="#242424")

style_lbl = {"fg": "white", "bg": "#242424", "font": ("Arial", 14, "bold")}
style_btn = {
    "bg": "#2a66a0", "fg": "white",
    "activebackground": "#1e4b75", "activeforeground": "white",
    "font": ("Arial", 11, "bold"), "width": 38, "padx": 10, "pady": 6
}

tk.Label(root, text="📝 PPTX & Excel Text Utility", **style_lbl).pack(pady=18)

tk.Button(root, text="Extract → All & Vietnamese‑only Excel", command=extract_text_both, **style_btn).pack(pady=4)
tk.Button(root, text="Apply 'TranslatedText' back to PPTX",  command=apply_translation,  **style_btn).pack(pady=4)
tk.Button(root, text="Copy 3rd‑col → replace 4th‑col in Excel", command=replace_excel_column, **style_btn).pack(pady=4)

root.mainloop()
